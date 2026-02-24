"""
Microserviço de renderização PPTX (estilo Manus).

Export template-first: abre um PPTX real, duplica o slide do template que corresponde
ao layout pedido e preenche o conteúdo. Animações e efeitos do template são preservados.

API:
- POST /render
  Body (JSON):
  {
    "templatePath": "/path/to/template.pptx",
    "deckTitle": "Título",
    "slides": [
      { "layout": "hero", "content": { "title": "...", "subtitle": "..." } },
      { "layout": "bullet", "content": { "title": "...", "items": [ { "text": "..." } ] } }
    ],
    "slideLayouts": [ { "index": 0, "type": "hero" }, { "index": 1, "type": "bullet" }, ... ]
  }

  slideLayouts = mapeamento no template: índice do slide → tipo de layout.
  Para cada slide em "slides", usamos o slide do template cujo type coincide com slide.layout.
  Se não houver correspondência, usa-se o slide de índice 0.
"""

from io import BytesIO
from typing import List, Optional

from fastapi import FastAPI, HTTPException
from fastapi.responses import Response
from pydantic import BaseModel
from pptx import Presentation


class SlideContent(BaseModel):
    layout: str
    content: dict


class SlideLayoutRef(BaseModel):
    index: int
    type: str


class RenderRequest(BaseModel):
    templatePath: str
    deckTitle: Optional[str] = None
    slides: List[SlideContent]
    slideLayouts: Optional[List[SlideLayoutRef]] = None


app = FastAPI(title="SlideDeck PPTX Renderer", version="1.0.0")


def _find_template_slide_index(layout: str, slide_layouts: List[dict]) -> int:
    """Devolve o índice do slide no template que corresponde ao layout pedido."""
    if not slide_layouts:
        return 0
    for ref in slide_layouts:
        ref_type = (ref.get("type") or ref.get("layout") or "").strip().lower()
        ref_type = ref_type.replace("_", "-")
        want = (layout or "title").strip().lower().replace("_", "-")
        if ref_type == want:
            idx = ref.get("index", 0)
            if isinstance(idx, int) and idx >= 0:
                return idx
    return 0


def _duplicate_slide(prs: Presentation, slide_index: int) -> "Slide":
    """
    Duplica um slide existente do template (preserva posição e estilo dos shapes).
    Copia caixas de texto; outros shapes podem ser ignorados nesta versão.
    """
    if slide_index >= len(prs.slides):
        slide_index = 0
    src = prs.slides[slide_index]
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    dest = prs.slides.add_slide(blank_layout)
    for shape in src.shapes:
        if not shape.has_text_frame:
            continue
        new_shape = dest.shapes.add_textbox(
            shape.left,
            shape.top,
            shape.width,
            shape.height,
        )
        new_tf = new_shape.text_frame
        new_tf.text = shape.text
    return dest


def _apply_content_to_slide(slide, layout: str, content: dict) -> None:
    """Preenche os shapes do slide com o conteúdo gerado pela IA."""
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    if not text_shapes:
        return

    def set_text(idx: int, value: str) -> None:
        if idx < len(text_shapes):
            tf = text_shapes[idx].text_frame
            tf.clear()
            if tf.paragraphs:
                tf.paragraphs[0].text = value or ""

    if layout in ("hero", "title", "title-subtitle", "title_subtitle"):
        set_text(0, content.get("title") or "")
        if layout != "title":
            set_text(1, content.get("subtitle") or "")
    elif layout == "bullet":
        set_text(0, content.get("title") or "")
        items = content.get("items") or []
        bullets = []
        for it in items:
            if isinstance(it, str):
                bullets.append(it)
            elif isinstance(it, dict):
                bullets.append(str(it.get("text") or ""))
        if len(text_shapes) > 1:
            tf = text_shapes[1].text_frame
            tf.clear()
            for line in bullets:
                p = tf.add_paragraph()
                p.text = line
    elif layout == "timeline":
        set_text(0, content.get("title") or "")
        events = content.get("events") or []
        lines = [f"{e.get('year', '')} – {e.get('text', '')}" if isinstance(e, dict) else str(e) for e in events]
        if len(text_shapes) > 1:
            tf = text_shapes[1].text_frame
            tf.clear()
            for line in lines:
                p = tf.add_paragraph()
                p.text = line
    elif layout in ("stats-row", "stats_row"):
        parts = [
            content.get("stat1") or "",
            content.get("label1") or "",
            content.get("stat2") or "",
            content.get("label2") or "",
            content.get("stat3") or "",
            content.get("label3") or "",
        ]
        set_text(0, "  |  ".join(p for p in parts if p))
    elif layout in ("big-number", "big_number"):
        set_text(0, str(content.get("number") or "0"))
        if len(text_shapes) > 1:
            set_text(1, content.get("label") or "")
    elif layout == "quote":
        set_text(0, content.get("text") or "")
        if len(text_shapes) > 1:
            set_text(1, content.get("author") or "")
    elif layout in ("section",):
        set_text(0, content.get("title") or "")
    elif layout == "two-column":
        left = content.get("left") or ""
        right = content.get("right") or ""
        set_text(0, left)
        if len(text_shapes) > 1:
            set_text(1, right)
    else:
        set_text(0, content.get("title") or "")


@app.post("/render")
def render_pptx(body: RenderRequest):
    try:
        prs = Presentation(body.templatePath)
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Erro ao abrir template PPTX: {exc}") from exc

    if body.deckTitle:
        prs.core_properties.title = body.deckTitle

    slide_layouts = body.slideLayouts or []
    layout_list = []
    for s in slide_layouts:
        idx = s.get("index", 0) if isinstance(s, dict) else getattr(s, "index", 0)
        typ = s.get("type", "") if isinstance(s, dict) else getattr(s, "type", "")
        layout_list.append({"index": idx, "type": typ})

    # Construir lista de índices do template a usar para cada slide pedido
    indices_to_use = []
    for s in body.slides:
        idx = _find_template_slide_index(s.layout, layout_list)
        indices_to_use.append(idx)

    # Remover todos os slides do template (vamos recriar a partir das cópias)
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    # Para cada slide pedido: duplicar o slide correto do template e preencher
    for i, s in enumerate(body.slides):
        # Reabrir o template para ler o slide de origem (já limpámos prs.slides)
        prs_ref = Presentation(body.templatePath)
        src_index = indices_to_use[i] if i < len(indices_to_use) else 0
        if src_index >= len(prs_ref.slides):
            src_index = 0
        src_slide = prs_ref.slides[src_index]
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        dest = prs.slides.add_slide(blank_layout)
        for shape in src_slide.shapes:
            if not shape.has_text_frame:
                continue
            new_shape = dest.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            new_shape.text_frame.text = shape.text
        _apply_content_to_slide(dest, s.layout, s.content or {})

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)

    return Response(
        content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="presentation.pptx"'},
    )
